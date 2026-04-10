"""
Génération automatique de la présentation DocuPost
Slides : Contexte | Agents (1/2) | Agents (2/2) | Mocks & Simulateurs
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Palette ──────────────────────────────────────────────────────────────────
C_BG        = RGBColor(0x0D, 0x1B, 0x2A)   # Navy profond
C_CARD      = RGBColor(0x1B, 0x2A, 0x3B)   # Bleu ardoise
C_ORANGE    = RGBColor(0xF4, 0x82, 0x1F)   # Accent primaire
C_TEAL      = RGBColor(0x00, 0xB4, 0xCC)   # Accent secondaire
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_GREY      = RGBColor(0xB0, 0xBE, 0xC5)
C_GREEN     = RGBColor(0x27, 0xAE, 0x60)
C_RED       = RGBColor(0xE7, 0x4C, 0x3C)
C_VIOLET    = RGBColor(0x9B, 0x59, 0xB6)
C_BLUE      = RGBColor(0x34, 0x98, 0xDB)
C_YELLOW    = RGBColor(0xF3, 0x9C, 0x12)
C_SLATE     = RGBColor(0x60, 0x7D, 0x8B)
C_ORANGE2   = RGBColor(0xE6, 0x7E, 0x22)

# ── Helpers ───────────────────────────────────────────────────────────────────

def add_bg(slide, color=C_BG):
    """Fond plein sur toute la slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=Pt(12), bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_text_multi(slide, lines, left, top, width, height,
                   font_size=Pt(11), color=C_WHITE, line_spacing=None):
    """lines = list of (text, bold, color_override)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for (text, bold, col) in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.size = font_size
        run.font.bold = bold
        run.font.color.rgb = col or color
    return txBox


def slide_title(slide, title, subtitle=None):
    """Bande header en haut."""
    add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.1), C_CARD)
    add_rect(slide, Inches(0), Inches(0.95), Inches(13.33), Inches(0.06), C_ORANGE)
    add_text(slide, title,
             Inches(0.3), Inches(0.1), Inches(10), Inches(0.6),
             font_size=Pt(26), bold=True, color=C_WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.3), Inches(0.62), Inches(10), Inches(0.38),
                 font_size=Pt(13), bold=False, color=C_TEAL)


def footer(slide):
    add_text(slide, "DocuPost MVP  •  Confidentiel  •  2026",
             Inches(0.3), Inches(7.2), Inches(12), Inches(0.3),
             font_size=Pt(9), color=C_GREY, align=PP_ALIGN.CENTER)


def card(slide, left, top, width, height, accent_color, title, lines, icon=""):
    """Carte agent avec barre colorée en haut."""
    add_rect(slide, left, top, width, height, C_CARD)
    add_rect(slide, left, top, width, Inches(0.08), accent_color)
    # Titre + icône
    add_text(slide, f"{icon}  {title}" if icon else title,
             left + Inches(0.12), top + Inches(0.12), width - Inches(0.2), Inches(0.38),
             font_size=Pt(13), bold=True, color=C_WHITE)
    # Corps
    y = top + Inches(0.52)
    for (label, content, col) in lines:
        add_text(slide, label,
                 left + Inches(0.12), y, width - Inches(0.2), Inches(0.22),
                 font_size=Pt(9), bold=True, color=col or C_ORANGE)
        y += Inches(0.22)
        add_text(slide, content,
                 left + Inches(0.12), y, width - Inches(0.2), Inches(0.5),
                 font_size=Pt(9), bold=False, color=C_WHITE)
        y += Inches(0.48)


# ── Présentation ──────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]  # blank

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Contexte, Cible & Contraintes
# ═══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
add_bg(s1)
slide_title(s1, "DocuPost — Contexte, Cibles & Contraintes",
            "Connecter le livreur terrain, le superviseur et le SI en temps réel")
footer(s1)

col_w  = Inches(4.2)
col_h  = Inches(5.8)
col_y  = Inches(1.2)
gap    = Inches(0.17)
col1_x = Inches(0.2)
col2_x = col1_x + col_w + gap
col3_x = col2_x + col_w + gap

# ── Colonne 1 : Contexte ──────────────────────────────────────────────────────
add_rect(s1, col1_x, col_y, col_w, col_h, C_CARD)
add_rect(s1, col1_x, col_y, Inches(0.06), col_h, C_ORANGE)

add_text(s1, "📋  CONTEXTE",
         col1_x + Inches(0.15), col_y + Inches(0.12), col_w, Inches(0.35),
         font_size=Pt(14), bold=True, color=C_ORANGE)

problemes = [
    "• Livreur hors SI dès le départ du dépôt → tournée sur papier",
    "• Supervision uniquement par téléphone, aucune visibilité temps réel",
    "• Preuves de livraison non disponibles immédiatement",
    "• Motifs de non-livraison hétérogènes et inexploitables",
    "• Double saisie manuelle dans les SI internes",
]
add_text(s1, "Problèmes identifiés :",
         col1_x + Inches(0.15), col_y + Inches(0.55), col_w - Inches(0.2), Inches(0.25),
         font_size=Pt(10), bold=True, color=C_TEAL)
add_text(s1, "\n".join(problemes),
         col1_x + Inches(0.15), col_y + Inches(0.82), col_w - Inches(0.2), Inches(1.8),
         font_size=Pt(9.5), color=C_WHITE)

add_text(s1, "Objectifs du projet :",
         col1_x + Inches(0.15), col_y + Inches(2.8), col_w - Inches(0.2), Inches(0.25),
         font_size=Pt(10), bold=True, color=C_TEAL)
objectifs = [
    "✅  Réintégrer le livreur dans le SI",
    "✅  Pilotage proactif (risque détecté < 15 min)",
    "✅  Preuves numériques opposables < 5 min",
    "✅  Zéro double saisie inter-systèmes",
    "✅  Sécurisation juridique des livraisons sensibles",
]
add_text(s1, "\n".join(objectifs),
         col1_x + Inches(0.15), col_y + Inches(3.1), col_w - Inches(0.2), Inches(2.0),
         font_size=Pt(9.5), color=C_WHITE)

# ── Colonne 2 : Utilisateurs cibles ──────────────────────────────────────────
add_rect(s1, col2_x, col_y, col_w, col_h, C_CARD)
add_rect(s1, col2_x, col_y, Inches(0.06), col_h, C_TEAL)

add_text(s1, "👥  UTILISATEURS CIBLES",
         col2_x + Inches(0.15), col_y + Inches(0.12), col_w, Inches(0.35),
         font_size=Pt(14), bold=True, color=C_TEAL)

users = [
    ("🚚  Pierre — Livreur terrain",
     "80–120 colis/jour · Android\nBesoin : liste claire, statut rapide, preuve sans friction"),
    ("📊  M. Renaud — Resp. Exploitation",
     "Pilote plusieurs tournées simultanément\nBesoin : vue temps réel, alertes proactives"),
    ("💼  Mme Dubois — DSI",
     "Audits, litiges, engagements contractuels\nBesoin : preuves opposables, SLA respectés"),
    ("🔧  M. Garnier — Architecte SI",
     "Garant cohérence SI · Valide intégrations\nBesoin : API REST, SSO, aucune modif OMS"),
]
uy = col_y + Inches(0.6)
for (title_u, desc) in users:
    add_rect(s1, col2_x + Inches(0.15), uy,
             col_w - Inches(0.3), Inches(1.18), RGBColor(0x0D, 0x1B, 0x2A))
    add_text(s1, title_u,
             col2_x + Inches(0.25), uy + Inches(0.06),
             col_w - Inches(0.5), Inches(0.3),
             font_size=Pt(10), bold=True, color=C_ORANGE)
    add_text(s1, desc,
             col2_x + Inches(0.25), uy + Inches(0.35),
             col_w - Inches(0.5), Inches(0.75),
             font_size=Pt(9), color=C_WHITE)
    uy += Inches(1.28)

# ── Colonne 3 : Contraintes ───────────────────────────────────────────────────
add_rect(s1, col3_x, col_y, col_w, col_h, C_CARD)
add_rect(s1, col3_x, col_y, Inches(0.06), col_h, C_VIOLET)

add_text(s1, "⚙️  CONTRAINTES",
         col3_x + Inches(0.15), col_y + Inches(0.12), col_w, Inches(0.35),
         font_size=Pt(14), bold=True, color=C_VIOLET)

contraintes_sections = [
    ("Stack technique imposée",
     "Java 21 / Spring Boot 4.0.3\nReact 19 / TypeScript 5.6\nDocker / Kubernetes"),
    ("Plateforme mobile",
     "Android uniquement (MVP)\niOS → Release 2"),
    ("Sécurité & Conformité",
     "OAuth2 / SSO Corporate obligatoire\nRGPD : géolocalisation minimisée\nTLS/HTTPS partout"),
    ("Intégration SI",
     "API REST OMS uniquement\nSans modifier le cœur OMS\nCRM / ERP → Release 2"),
    ("Environnements",
     "dev → recette → préprod → prod\n(4 environnements imposés DSI)"),
]
cy = col_y + Inches(0.56)
for (sec_title, sec_content) in contraintes_sections:
    add_text(s1, sec_title,
             col3_x + Inches(0.15), cy, col_w - Inches(0.2), Inches(0.25),
             font_size=Pt(10), bold=True, color=C_TEAL)
    cy += Inches(0.25)
    add_text(s1, sec_content,
             col3_x + Inches(0.15), cy, col_w - Inches(0.2), Inches(0.55),
             font_size=Pt(9), color=C_WHITE)
    cy += Inches(0.6)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agents : Vision → Architecture → Backlog (Sponsor, UX, Archi M, Archi T, PO)
# ═══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
add_bg(s2)
slide_title(s2, "Chaîne de valeur — Agents 1/2 : Vision → Architecture → Backlog",
            "Sponsor · UX · Architecte Métier · Architecte Technique · Product Owner")
footer(s2)

agents_s2 = [
    {
        "icon": "🎯", "title": "SPONSOR", "color": C_ORANGE,
        "role": "Garant vision business, KPIs et périmètre MVP. Arbitre entre Core Domain et sous-domaines.",
        "inputs": "Entretiens métier terrain\n(Pierre, M. Renaud, Mme Dubois, M. Garnier)",
        "outputs": "vision-produit.md\nkpis.md\nperimetre-mvp.md",
    },
    {
        "icon": "🎨", "title": "UX DESIGNER", "color": C_TEAL,
        "role": "Conçoit l'expérience livreur et superviseur. Capture le langage terrain (Ubiquitous Language).",
        "inputs": "Vision produit\nPain points entretiens\nParcours utilisateurs",
        "outputs": "personas.md\nuser-journeys.md\nwireframes.md\ndesign-system.md",
    },
    {
        "icon": "🗺️", "title": "ARCHI. MÉTIER", "color": C_VIOLET,
        "role": "Structure les domaines, entités et règles métier (DDD). Construit la capability map.",
        "inputs": "Vision produit\nUser journeys UX\nLanguage terrain",
        "outputs": "domain-model.md\ncapability-map.md\nmodules-fonctionnels.md",
    },
    {
        "icon": "⚙️", "title": "ARCHI. TECHNIQUE", "color": C_BLUE,
        "role": "Définit l'architecture applicative, les intégrations et choix technologiques (Bounded Contexts → services).",
        "inputs": "UX + Archi Métier\nContraintes SI\n(SSO, OMS, stack)",
        "outputs": "architecture-applicative.md\nschemas-integration.md\ndesign-decisions.md\nNFR.md",
    },
    {
        "icon": "📋", "title": "PRODUCT OWNER", "color": C_YELLOW,
        "role": "Transforme la vision en backlog structuré Epics / Features / User Stories (SAFe, MoSCoW).",
        "inputs": "Vision + UX\nArchi Métier\nGlossaire domaine",
        "outputs": "epics.md · features.md\nUS-[NNN]-[slug].md\ndefinition-mvp.md",
    },
]

card_w = Inches(2.46)
card_h = Inches(5.8)
card_y = Inches(1.2)
card_gap = Inches(0.12)

for i, ag in enumerate(agents_s2):
    cx = Inches(0.2) + i * (card_w + card_gap)
    add_rect(s2, cx, card_y, card_w, card_h, C_CARD)
    add_rect(s2, cx, card_y, card_w, Inches(0.1), ag["color"])

    # Icon + Titre
    add_text(s2, ag["icon"],
             cx + Inches(0.1), card_y + Inches(0.14), card_w - Inches(0.15), Inches(0.38),
             font_size=Pt(22), bold=False, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s2, ag["title"],
             cx + Inches(0.1), card_y + Inches(0.55), card_w - Inches(0.15), Inches(0.32),
             font_size=Pt(12), bold=True, color=ag["color"], align=PP_ALIGN.CENTER)

    # Rôle
    add_text(s2, "RÔLE",
             cx + Inches(0.12), card_y + Inches(0.95), card_w - Inches(0.2), Inches(0.22),
             font_size=Pt(8.5), bold=True, color=C_GREY)
    add_text(s2, ag["role"],
             cx + Inches(0.12), card_y + Inches(1.18), card_w - Inches(0.2), Inches(1.1),
             font_size=Pt(9), color=C_WHITE)

    # Inputs
    add_rect(s2, cx + Inches(0.1), card_y + Inches(2.4),
             card_w - Inches(0.2), Inches(0.04), ag["color"])
    add_text(s2, "INPUTS",
             cx + Inches(0.12), card_y + Inches(2.5), card_w - Inches(0.2), Inches(0.22),
             font_size=Pt(8.5), bold=True, color=C_GREY)
    add_text(s2, ag["inputs"],
             cx + Inches(0.12), card_y + Inches(2.74), card_w - Inches(0.2), Inches(0.9),
             font_size=Pt(9), color=C_WHITE)

    # Outputs
    add_rect(s2, cx + Inches(0.1), card_y + Inches(3.75),
             card_w - Inches(0.2), Inches(0.04), ag["color"])
    add_text(s2, "OUTPUTS",
             cx + Inches(0.12), card_y + Inches(3.85), card_w - Inches(0.2), Inches(0.22),
             font_size=Pt(8.5), bold=True, color=C_GREY)
    add_text(s2, ag["outputs"],
             cx + Inches(0.12), card_y + Inches(4.09), card_w - Inches(0.2), Inches(1.4),
             font_size=Pt(9), color=C_WHITE)

# Flèche de flux en bas
add_text(s2, "Sponsor  →  UX  ┐                                    ┌→  PO  →  Dev  →  QA  →  DevOps  →  End User",
         Inches(0.5), Inches(7.05), Inches(12.5), Inches(0.25),
         font_size=Pt(8.5), color=C_GREY, italic=True)
add_text(s2, "                          └→  Archi Technique  ──────┘",
         Inches(0.5), Inches(7.22), Inches(12.5), Inches(0.2),
         font_size=Pt(8.5), color=C_GREY, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Agents : Dev → QA → DevOps → End User
# ═══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
add_bg(s3)
slide_title(s3, "Chaîne de valeur — Agents 2/2 : Développement → Test → Production",
            "Développeur · QA · DevOps · End User  —  Boucle itérative par User Story")
footer(s3)

agents_s3 = [
    {
        "icon": "💻", "title": "DÉVELOPPEUR", "color": C_GREEN,
        "role": "Implémente 1 User Story de bout en bout (vertical slice) : backend Java + frontend React/Expo, avec tests unitaires.",
        "inputs": "US-[NNN]-[slug].md\nWireframes UX\nArchitecture technique\nNFR",
        "outputs": "Code backend + frontend\nUS-[NNN]-impl.md\n(commandes démarrage, URLs test)",
    },
    {
        "icon": "🧪", "title": "QA", "color": C_RED,
        "role": "Définit et exécute la stratégie de tests (pyramide L1/L2/L3). Rédige scénarios Gherkin, exécute Playwright E2E.",
        "inputs": "US-[NNN]-[slug].md\nUS-[NNN]-impl.md\nExigences non fonctionnelles",
        "outputs": "plan-tests.md\nUS-[NNN]-scenarios.md\nRapport Playwright\nCheck-list tests manuels",
    },
    {
        "icon": "🚀", "title": "DEVOPS", "color": C_SLATE,
        "role": "Conçoit pipeline CI/CD GitHub Actions, provisionne les environnements GCP Cloud Run, définit le monitoring.",
        "inputs": "Architecture technique\nBounded Contexts\nNFR (SLA, disponibilité)",
        "outputs": "pipeline-cicd.md\nstrategie-deploiement.md\nmonitoring.md\ncloudbuild.yaml",
    },
    {
        "icon": "👷", "title": "END USER", "color": C_ORANGE2,
        "role": "Joue le rôle livreur ou superviseur en conditions réelles. Valide ergonomie, détecte frictions et incohérences de langage.",
        "inputs": "Wireframes UX\nUS implémentée\nScénarios QA",
        "outputs": "feedback-[feature]-[date].md\n→ Bloquants\n→ Améliorations\n→ Mineurs",
    },
]

card_w3 = Inches(3.1)
card_h3 = Inches(4.5)
card_y3 = Inches(1.25)
card_gap3 = Inches(0.14)

for i, ag in enumerate(agents_s3):
    cx = Inches(0.2) + i * (card_w3 + card_gap3)
    add_rect(s3, cx, card_y3, card_w3, card_h3, C_CARD)
    add_rect(s3, cx, card_y3, card_w3, Inches(0.1), ag["color"])

    add_text(s3, ag["icon"],
             cx + Inches(0.1), card_y3 + Inches(0.14), card_w3 - Inches(0.15), Inches(0.38),
             font_size=Pt(26), color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s3, ag["title"],
             cx + Inches(0.1), card_y3 + Inches(0.58), card_w3 - Inches(0.15), Inches(0.32),
             font_size=Pt(13), bold=True, color=ag["color"], align=PP_ALIGN.CENTER)

    add_text(s3, "RÔLE",
             cx + Inches(0.15), card_y3 + Inches(1.0), card_w3 - Inches(0.25), Inches(0.22),
             font_size=Pt(8.5), bold=True, color=C_GREY)
    add_text(s3, ag["role"],
             cx + Inches(0.15), card_y3 + Inches(1.22), card_w3 - Inches(0.25), Inches(1.0),
             font_size=Pt(9.5), color=C_WHITE)

    add_rect(s3, cx + Inches(0.12), card_y3 + Inches(2.3),
             card_w3 - Inches(0.24), Inches(0.04), ag["color"])
    add_text(s3, "INPUTS",
             cx + Inches(0.15), card_y3 + Inches(2.4), card_w3 - Inches(0.25), Inches(0.22),
             font_size=Pt(8.5), bold=True, color=C_GREY)
    add_text(s3, ag["inputs"],
             cx + Inches(0.15), card_y3 + Inches(2.62), card_w3 - Inches(0.25), Inches(0.8),
             font_size=Pt(9), color=C_WHITE)

    add_rect(s3, cx + Inches(0.12), card_y3 + Inches(3.5),
             card_w3 - Inches(0.24), Inches(0.04), ag["color"])
    add_text(s3, "OUTPUTS",
             cx + Inches(0.15), card_y3 + Inches(3.6), card_w3 - Inches(0.25), Inches(0.22),
             font_size=Pt(8.5), bold=True, color=C_GREY)
    add_text(s3, ag["outputs"],
             cx + Inches(0.15), card_y3 + Inches(3.82), card_w3 - Inches(0.25), Inches(0.6),
             font_size=Pt(9), color=C_WHITE)

# Boucle itérative
loop_y = Inches(5.88)
add_rect(s3, Inches(0.2), loop_y, Inches(12.93), Inches(1.3), RGBColor(0x0D, 0x28, 0x40))
add_rect(s3, Inches(0.2), loop_y, Inches(0.08), Inches(1.3), C_TEAL)
add_text(s3, "BOUCLE ITÉRATIVE PAR USER STORY",
         Inches(0.45), loop_y + Inches(0.08), Inches(12.5), Inches(0.28),
         font_size=Pt(11), bold=True, color=C_TEAL)
add_text(s3,
         "PO définit l'US  →  Dev implémente  →  QA teste (Playwright)\n"
         "  Si FAIL → retour Dev (corrections)       Si PASS → End User valide  →  PO crée la prochaine US",
         Inches(0.45), loop_y + Inches(0.38), Inches(12.5), Inches(0.85),
         font_size=Pt(10), color=C_WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Mocks & Simulateurs
# ═══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
add_bg(s4)
slide_title(s4, "Mocks & Simulateurs — Infrastructure dev/test",
            "@Profile(\"dev\") — Aucun simulateur actif en production  •  404 Not Found sur /dev/** en prod")
footer(s4)

# ── Zone gauche : Inventaire des mocks ────────────────────────────────────────
left_x = Inches(0.2)
left_w = Inches(5.8)
zone_y = Inches(1.25)
zone_h = Inches(5.85)

add_rect(s4, left_x, zone_y, left_w, zone_h, C_CARD)
add_rect(s4, left_x, zone_y, Inches(0.07), zone_h, C_ORANGE)

add_text(s4, "MOCKS EN PLACE",
         left_x + Inches(0.18), zone_y + Inches(0.1), left_w - Inches(0.25), Inches(0.3),
         font_size=Pt(13), bold=True, color=C_ORANGE)

mocks = [
    ("[1] Simulateur TMS  —  POST /dev/tms/import", C_ORANGE,
     "Prob : en prod, les tournees viennent d'un logiciel TMS externe\n"
     "       (non accessible en dev).\n"
     "Mock : cet endpoint cree de fausses tournees directement en base\n"
     "       (zones Lyon, 3-8 colis, contraintes horaires realistes)."),
    ("[2] DevEventBridge  —  @Profile(\"dev\")", C_TEAL,
     "Prob : en prod, les 3 services communiquent via Kafka.\n"
     "       Kafka n'est pas deploye en dev/recette.\n"
     "Mock : quand une tournee est lancee, ce composant appelle\n"
     "       directement svc-tournee et svc-supervision par HTTP."),
    ("[3] Endpoints Reset  —  DELETE & POST /dev/tms/reset", C_BLUE,
     "Prob : apres chaque campagne de tests Playwright, la base\n"
     "       contient des residus qui faussent les tests suivants.\n"
     "Mock : ces endpoints purgent toutes les donnees de test en\n"
     "       une seule requete pour repartir d'un etat propre."),
    ("[4] Reseed svc-tournee  —  POST /internal/dev/reseed", C_GREEN,
     "Prob : le service livreur (svc-tournee) n'a aucune tournee\n"
     "       a afficher apres un reset.\n"
     "Mock : cet endpoint reinsere des tournees fictives avec des\n"
     "       colis realistes pour que le livreur puisse tester."),
    ("[5] OMS simule  —  donnees fictives en base", C_VIOLET,
     "Prob : en prod, chaque livraison envoie un evenement au\n"
     "       systeme OMS Docaposte (externe, inaccessible en dev).\n"
     "Mock : les statuts colis sont stockes localement sans appel\n"
     "       OMS reel. L'interface se comporte comme si tout partait."),
    ("[6] Auth bypass  —  REACT_APP_AUTH_BYPASS=true", C_YELLOW,
     "Prob : en prod, la connexion passe par le SSO corporate\n"
     "       Docaposte (inaccessible depuis le poste dev).\n"
     "Mock : cette variable desactive le SSO et connecte l'utilisateur\n"
     "       automatiquement pour tester sans compte SI reel."),
]

my = zone_y + Inches(0.48)
for (label, col, desc) in mocks:
    add_text(s4, label,
             left_x + Inches(0.18), my, left_w - Inches(0.28), Inches(0.24),
             font_size=Pt(10), bold=True, color=col)
    my += Inches(0.24)
    add_text(s4, desc,
             left_x + Inches(0.18), my, left_w - Inches(0.28), Inches(0.6),
             font_size=Pt(8.5), color=C_WHITE)
    my += Inches(0.66)

# ── Zone droite : Schéma flux ──────────────────────────────────────────────────
right_x = left_x + left_w + Inches(0.2)
right_w = Inches(6.93)

# Bloc PROD
add_rect(s4, right_x, zone_y, right_w, Inches(1.5), RGBColor(0x28, 0x18, 0x18))
add_rect(s4, right_x, zone_y, right_w, Inches(0.07), C_RED)
add_text(s4, "PROD — Flux réel (non disponible en dev)",
         right_x + Inches(0.15), zone_y + Inches(0.1), right_w - Inches(0.2), Inches(0.28),
         font_size=Pt(10), bold=True, color=C_RED)
add_text(s4,
         "TMS externe  →  Kafka  →  BC-07 (import TourneePlanifiee)\n"
         "                                   →  BC-01 (Tournee mobile)  +  BC-03 (VueTournee)",
         right_x + Inches(0.15), zone_y + Inches(0.42), right_w - Inches(0.2), Inches(0.9),
         font_size=Pt(9.5), color=C_WHITE)

# Flèche vers DEV
add_text(s4, "▼  REMPLACÉ PAR LES SIMULATEURS",
         right_x + Inches(0.15), zone_y + Inches(1.6), right_w - Inches(0.2), Inches(0.3),
         font_size=Pt(10), bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

# Bloc DEV
add_rect(s4, right_x, zone_y + Inches(1.98), right_w, Inches(4.12), RGBColor(0x0D, 0x28, 0x18))
add_rect(s4, right_x, zone_y + Inches(1.98), right_w, Inches(0.07), C_GREEN)
add_text(s4, "DEV / RECETTE — Flux simulé (@Profile(\"dev\"))",
         right_x + Inches(0.15), zone_y + Inches(2.08), right_w - Inches(0.2), Inches(0.28),
         font_size=Pt(10), bold=True, color=C_GREEN)

schema = (
    "POST /dev/tms/import\n"
    "         ↓\n"
    "TourneePlanifiee créée (BC-07, svc-supervision)\n"
    "         ↓\n"
    "LancerTourneeHandler  +  DevEventBridge (HTTP, @Profile dev)\n"
    "         ↓                              ↓\n"
    "POST /internal/dev/tournees    VueTournee créée (BC-03)\n"
    "Tournee créée (BC-01)          (idempotent)\n"
    "         ↓\n"
    "Livreur voit sa tournée        Superviseur voit le tableau de bord\n"
    "(svc-tournee:8081)             (svc-supervision:8082)"
)
add_text(s4, schema,
         right_x + Inches(0.15), zone_y + Inches(2.42), right_w - Inches(0.2), Inches(3.5),
         font_size=Pt(9.5), color=C_WHITE)

# Règle sécurité
add_rect(s4, right_x, zone_y + Inches(4.3), right_w, Inches(1.55),
         RGBColor(0x28, 0x10, 0x10))
add_rect(s4, right_x, zone_y + Inches(4.3), right_w, Inches(0.06), C_RED)
add_text(s4, "⚠️  RÈGLES DE SÉCURITÉ",
         right_x + Inches(0.15), zone_y + Inches(4.4), right_w - Inches(0.2), Inches(0.28),
         font_size=Pt(10), bold=True, color=C_RED)
add_text(s4,
         "@Profile(\"dev\") obligatoire sur TOUS les beans et endpoints de simulation\n"
         "404 Not Found en prod sur tous les /dev/**  •  Aucun log de simulation en prod\n"
         "Idempotence garantie sur tous les endpoints dev (double appel → 200 OK sans doublon)",
         right_x + Inches(0.15), zone_y + Inches(4.72), right_w - Inches(0.2), Inches(1.0),
         font_size=Pt(9), color=C_WHITE)


# ── Sauvegarde ────────────────────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "DocuPost_Presentation.pptx")
prs.save(output_path)
print(f"OK  Presentation generee : {output_path}")
